import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage
import os
import shutil
import subprocess
import tempfile
from html2image import Html2Image


def _find_chrome():
    """Find Chrome/Chromium executable, installing in Colab if needed."""
    candidates = [
        '/usr/bin/google-chrome-stable',
        '/usr/bin/google-chrome',
        '/usr/bin/chromium-browser',
        '/usr/bin/chromium',
        shutil.which('google-chrome-stable'),
        shutil.which('google-chrome'),
        shutil.which('chromium-browser'),
        shutil.which('chromium'),
    ]
    for path in candidates:
        if path and os.path.exists(path):
            return path

    # Auto-install in Colab
    try:
        import google.colab  # noqa: F401
        subprocess.run(
            ['bash', '-c',
             'wget -q -O /tmp/chrome.deb '
             'https://dl.google.com/linux/direct/google-chrome-stable_current_amd64.deb '
             '&& apt-get install -y -qq /tmp/chrome.deb '
             '&& rm /tmp/chrome.deb'],
            check=True, capture_output=True,
        )
        if os.path.exists('/usr/bin/google-chrome-stable'):
            return '/usr/bin/google-chrome-stable'
    except (ImportError, subprocess.CalledProcessError):
        pass

    raise FileNotFoundError(
        "No Chrome/Chromium found. In Colab, google-chrome-stable should be "
        "pre-installed. Try: !which google-chrome-stable"
    )


# ──────────────────────────────────────────────────────────────────────
# SkillCorner brand colours — exact values from template
# ──────────────────────────────────────────────────────────────────────
SC_BG = RGBColor(0x25, 0x25, 0x25)          # #252525 — title bar / dark panels
SC_GREEN = RGBColor(0x32, 0xFE, 0x6B)       # #32FE6B — neon green accent
SC_GREEN_LINE = RGBColor(0x32, 0xFF, 0x6A)  # #32FF6A — accent lines
SC_DARK_GREEN = RGBColor(0x00, 0xA8, 0x2F)  # #00A82F
SC_PANEL_GREY = RGBColor(0x88, 0x88, 0x88)  # #888888 — content panels
SC_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SC_GREY = RGBColor(0x87, 0x87, 0x87)        # #878787
SC_LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xE8)  # #E8E8E8
SC_DARK_TEXT = RGBColor(0x57, 0x57, 0x57)    # #575757

# Template slide dimensions (10 x 5.625 inches — standard 16:9)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Template measurements (from inspecting the PPTX)
TITLE_LEFT = Inches(0.326)       # Title bar left edge
TITLE_TOP = Inches(0.082)        # Title bar top
TITLE_HEIGHT = Inches(0.572)     # Title bar height
SUBTITLE_TOP = Inches(0.503)     # Subtitle top
SUBTITLE_HEIGHT = Inches(0.283)  # Subtitle height
CONTENT_TOP = Inches(0.875)      # Where content area begins

# ──────────────────────────────────────────────────────────────────────
# Branded asset paths (extracted from template)
# ──────────────────────────────────────────────────────────────────────
try:
    _ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               'template_assets')
except NameError:
    # __file__ not defined (e.g. interactive/Colab) — try cwd
    _ASSETS_DIR = os.path.join(os.getcwd(), 'template_assets')

# Dark bg with green glow — cover slides
BG_COVER = os.path.join(_ASSETS_DIR, 'slide1_bg_6223517a.png')
# Dark bg with geometric diamonds + green glow — section dividers
BG_SECTION = os.path.join(_ASSETS_DIR, 'slide9_rel_e6607ea9.png')
# SKILLCORNER wordmark (white, transparent)
LOGO_WORDMARK = os.path.join(_ASSETS_DIR, 'slide6_logo_d2b96a8f.png')
# SC icon logo (white, transparent)
LOGO_ICON = os.path.join(_ASSETS_DIR, 'slide6_logo_019604da.png')
# Content slide background — white with SKILLCORNER wordmark + lines
BG_CONTENT = os.path.join(_ASSETS_DIR, 'slide53_embed_rId3_55e09a58.png')


def _strip_download_buttons(html_string):
    """Remove download/export buttons from HTML before rendering."""
    html_string = re.sub(
        r'<button[^>]*class="[^"]*download[^"]*"[^>]*>.*?</button>',
        '', html_string, flags=re.IGNORECASE | re.DOTALL
    )
    html_string = re.sub(
        r'<button[^>]*onclick="[^"]*download[^"]*"[^>]*>.*?</button>',
        '', html_string, flags=re.IGNORECASE | re.DOTALL
    )
    html_string = re.sub(
        r'<button[^>]*>.*?(?:download|export|save\s+png).*?</button>',
        '', html_string, flags=re.IGNORECASE | re.DOTALL
    )
    html_string = re.sub(
        r'<a[^>]*download[^>]*>.*?</a>',
        '', html_string, flags=re.IGNORECASE | re.DOTALL
    )
    html_string = re.sub(
        r'\.download-btn\s*\{[^}]*\}',
        '', html_string, flags=re.IGNORECASE
    )
    html_string = re.sub(
        r'\.download-button\s*\{[^}]*\}',
        '', html_string, flags=re.IGNORECASE
    )
    html_string = re.sub(
        r'\.pitch-download\s*\{[^}]*\}',
        '', html_string, flags=re.IGNORECASE
    )
    html_string = re.sub(
        r'\.download-button::before\s*\{[^}]*\}',
        '', html_string, flags=re.IGNORECASE
    )
    html_string = re.sub(
        r'\.download-button:hover\s*\{[^}]*\}',
        '', html_string, flags=re.IGNORECASE
    )
    html_string = re.sub(
        r'[^{]*\.pitch-download[^{]*\{[^}]*\}',
        '', html_string, flags=re.IGNORECASE
    )
    html_string = re.sub(
        r'function\s+downloadPNG\s*\(\)\s*\{[^}]*(?:\{[^}]*\}[^}]*)*\}',
        '', html_string, flags=re.IGNORECASE
    )
    return html_string


def _html_to_image(html_string, output_path, size=(1400, 900)):
    """Render an HTML string to a PNG image using Chrome headless."""
    chrome_path = _find_chrome()
    output_dir = os.path.dirname(os.path.abspath(output_path))
    output_name = os.path.basename(output_path)

    html_string = _strip_download_buttons(html_string)

    inject_css = (
        '<link href="https://fonts.googleapis.com/css2?family=Chakra+Petch:'
        'wght@300;400;500;600;700&display=swap" rel="stylesheet">'
        '<style>'
        'html, body { margin: 0; padding: 0; background: #ffffff; } '
        '::-webkit-scrollbar { display: none !important; '
        'width: 0 !important; height: 0 !important; } '
        'body { font-family: "Chakra Petch", sans-serif; '
        '-ms-overflow-style: none; scrollbar-width: none; } '
        '</style>'
    )
    if '</head>' in html_string:
        html_string = html_string.replace('</head>', inject_css + '</head>')
    else:
        html_string = inject_css + html_string

    hti = Html2Image(
        browser_executable=chrome_path,
        output_path=output_dir,
        size=size,
        custom_flags=[
            '--no-sandbox',
            '--disable-gpu',
            '--hide-scrollbars',
            '--force-device-scale-factor=2',
            '--virtual-time-budget=5000',
        ],
    )
    hti.screenshot(html_str=html_string, save_as=output_name)

    _trim_whitespace(output_path)
    return output_path


def _trim_whitespace(image_path):
    """Crop trailing white/transparent space from bottom and right.
    Uses a conservative threshold — only trims rows/columns that are
    entirely white (RGB > 252 on all channels) or fully transparent."""
    import numpy as np

    with PILImage.open(image_path) as img:
        img = img.convert('RGBA')
        arr = np.array(img)
        h, w = arr.shape[:2]

        # A pixel is "empty" if fully transparent OR near-white
        r, g, b, a = arr[:, :, 0], arr[:, :, 1], arr[:, :, 2], arr[:, :, 3]
        has_content = (a > 10) & ~((r > 252) & (g > 252) & (b > 252))

        # Find last row/col with any content
        row_has_content = has_content.any(axis=1)
        col_has_content = has_content.any(axis=0)

        if not row_has_content.any():
            return  # entirely blank image, don't touch

        last_row = int(np.where(row_has_content)[0][-1])
        last_col = int(np.where(col_has_content)[0][-1])

        # Add small margin
        bottom = min(last_row + 6, h)
        right = min(last_col + 6, w)

        if bottom < h - 20 or right < w - 20:
            cropped = img.crop((0, 0, right, bottom))
            cropped.save(image_path)


# ──────────────────────────────────────────────────────────────────────
# Slide helpers — directly replicating SkillCorner template
# ──────────────────────────────────────────────────────────────────────

def _get_blank_layout(prs):
    """Find the most blank layout — prefer 0 placeholders, then by name."""
    # First try to find a layout with 0 placeholders
    for layout in prs.slide_layouts:
        if len(list(layout.placeholders)) == 0:
            return layout

    # Then try by name
    for name in ['BLANK', 'Blank']:
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout

    # Fallback: layout with fewest placeholders
    return min(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))


def _add_filled_rect(slide, left, top, width, height, fill_color):
    """Add a solid-filled rectangle with no border and no shadow."""
    from pptx.oxml.ns import qn
    from lxml import etree
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.shadow.inherit = False

    spPr = shape._element.spPr

    # Remove border completely via XML
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'), w='0')
    etree.SubElement(ln, qn('a:noFill'))

    # Remove any effect list (shadows, glows, etc.)
    for tag in [qn('a:effectLst'), qn('a:effectDag')]:
        el = spPr.find(tag)
        if el is not None:
            spPr.remove(el)
    # Add empty effectLst to explicitly disable all effects
    etree.SubElement(spPr, qn('a:effectLst'))

    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  font_color=SC_WHITE, bold=True, alignment=PP_ALIGN.LEFT,
                  font_name='Chakra Petch'):
    """Add a text box with font set at paragraph, run, AND XML level."""
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment

    # Set font at paragraph default run properties (defRPr)
    pPr = p._p.get_or_add_pPr()
    defRPr = pPr.find(qn('a:defRPr'))
    if defRPr is None:
        defRPr = pPr.makeelement(qn('a:defRPr'), {})
        pPr.append(defRPr)
    # Set latin font on defRPr
    latin = defRPr.find(qn('a:latin'))
    if latin is None:
        latin = defRPr.makeelement(qn('a:latin'), {})
        defRPr.append(latin)
    latin.set('typeface', font_name)

    # Set font at run level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name

    # Also set latin typeface directly on the run XML for Google Slides
    rPr = run._r.get_or_add_rPr()
    r_latin = rPr.find(qn('a:latin'))
    if r_latin is None:
        r_latin = rPr.makeelement(qn('a:latin'), {})
        rPr.append(r_latin)
    r_latin.set('typeface', font_name)

    return txBox


def _set_slide_bg(slide, color):
    """Set the slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_bg_image(slide, image_path):
    """Add a full-slide background image (sent to back)."""
    if image_path and os.path.exists(image_path):
        pic = slide.shapes.add_picture(
            image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
        )
        # Move to back of z-order
        sp = pic._element
        slide.shapes._spTree.remove(sp)
        slide.shapes._spTree.insert(2, sp)


def _add_logo(slide, logo_path, left, top, width=None, height=None):
    """Add a logo image at specified position."""
    if logo_path and os.path.exists(logo_path):
        kwargs = {}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(logo_path, left, top, **kwargs)


def _add_title_slide(prs, team_name, report_title=None, report_subtitle=None):
    """
    Cover slide — branded dark background with green glow,
    SC wordmark logo, matching template cover slides 6-8.
    """
    slide = prs.slides.add_slide(_get_blank_layout(prs))  # BLANK layout
    _set_slide_bg(slide, SC_BG)

    # Branded background image (dark with green glow)
    _add_bg_image(slide, BG_COVER)

    # Title — Chakra Petch Bold ~41pt
    title_text = report_title if report_title else f"{team_name} Analysis Report"
    _add_text_box(slide, Inches(0.6), Inches(1.6), Inches(8.5), Inches(1.2),
                  title_text, font_size=41, font_color=SC_WHITE, bold=True)

    # Subtitle — Chakra Petch, #32FE6B green
    subtitle_text = report_subtitle if report_subtitle else team_name
    _add_text_box(slide, Inches(0.6), Inches(2.8), Inches(8.5), Inches(0.8),
                  subtitle_text, font_size=16, font_color=SC_GREEN, bold=False)

    # SC wordmark logo — bottom left (template position)
    _add_logo(slide, LOGO_WORDMARK,
              Inches(0.6), SLIDE_HEIGHT - Inches(0.45),
              width=Inches(1.17))

    # SC icon logo — bottom right
    _add_logo(slide, LOGO_ICON,
              SLIDE_WIDTH - Inches(0.6) - Inches(0.22),
              SLIDE_HEIGHT - Inches(0.5),
              width=Inches(0.22))


def _add_section_divider(prs, section_title, section_number=None):
    """
    Section divider — branded dark background with geometric pattern,
    large green number + white title. Matches template slide 11 style.
    """
    slide = prs.slides.add_slide(_get_blank_layout(prs))  # BLANK layout
    _set_slide_bg(slide, SC_BG)

    # Branded background image (geometric diamonds + green glow)
    _add_bg_image(slide, BG_SECTION)

    # Large green section number — Chakra Petch Bold ~72pt
    if section_number is not None:
        _add_text_box(slide, Inches(0.6), Inches(0.8), Inches(3), Inches(1.8),
                      str(section_number).zfill(2), font_size=72,
                      font_color=SC_GREEN, bold=True)

    # Section title — Chakra Petch Bold ~34pt
    title_top = Inches(1.2) if section_number is None else Inches(2.6)
    _add_text_box(slide, Inches(0.6), title_top, Inches(8.5), Inches(1.2),
                  section_title, font_size=34, font_color=SC_WHITE, bold=True)

    # SC icon logo — bottom right
    _add_logo(slide, LOGO_ICON,
              SLIDE_WIDTH - Inches(0.6) - Inches(0.22),
              SLIDE_HEIGHT - Inches(0.5),
              width=Inches(0.22))


def _add_content_slide(prs, title, image_path, subtitle=None):
    """
    Content / analysis slide — white background with SKILLCORNER
    branded overlay (wordmark + lines from template background PNG).
    Title area has white fill on top of the overlay so text is readable.
    """
    slide = prs.slides.add_slide(_get_blank_layout(prs))

    # White slide background
    _set_slide_bg(slide, SC_WHITE)

    # Branded background overlay — SKILLCORNER wordmark + lines
    _add_bg_image(slide, BG_CONTENT)

    # White fill behind title area (on top of BG overlay, under text)
    _add_filled_rect(slide, 0, 0, Inches(10.0 / 3), Inches(0.65), SC_WHITE)

    # Title — Chakra Petch Bold, black (#252525)
    _add_text_box(slide, Inches(0.35), Inches(0.12), Inches(6), Inches(0.4),
                  title, font_size=20, font_color=SC_BG, bold=True)

    # Subtitle — Chakra Petch Regular, grey
    if subtitle:
        _add_text_box(slide, Inches(0.35), Inches(0.45), Inches(6), Inches(0.2),
                      subtitle, font_size=10, font_color=SC_GREY, bold=False)

    # Image — maximise available space on slide
    if image_path and os.path.exists(image_path):
        with PILImage.open(image_path) as img:
            img_w, img_h = img.size

        # Content area: padded to keep SKILLCORNER border line visible
        content_top_in = 0.65
        pad_left = 0.25
        pad_right = 0.45   # right side has vertical border line
        pad_bottom = 0.1
        max_w = 10.0 - pad_left - pad_right
        max_h = 5.625 - content_top_in - pad_bottom

        img_aspect = img_w / img_h
        area_aspect = max_w / max_h

        if img_aspect > area_aspect:
            fit_w = max_w
            fit_h = max_w / img_aspect
        else:
            fit_h = max_h
            fit_w = max_h * img_aspect

        # Centre horizontally and vertically in content area
        left = pad_left + (max_w - fit_w) / 2
        top = content_top_in + (max_h - fit_h) / 2

        slide.shapes.add_picture(
            image_path,
            Inches(left), Inches(top),
            width=Inches(fit_w), height=Inches(fit_h),
        )


def _add_dual_image_slide(prs, title, left_image_path, right_image_path,
                          left_label=None, right_label=None):
    """
    Content slide with two images side by side (e.g. IP + OOP radars).
    Each image gets a label above it.
    """
    slide = prs.slides.add_slide(_get_blank_layout(prs))
    _set_slide_bg(slide, SC_WHITE)
    _add_bg_image(slide, BG_CONTENT)

    # White fill behind title
    _add_filled_rect(slide, 0, 0, Inches(10.0 / 3), Inches(0.65), SC_WHITE)

    # Title
    _add_text_box(slide, Inches(0.35), Inches(0.12), Inches(6), Inches(0.4),
                  title, font_size=20, font_color=SC_BG, bold=True)

    # Layout: two images side by side
    content_top = 0.65
    pad_side = 0.3
    gap = 0.2
    available_w = 10.0 - pad_side * 2 - gap
    half_w = available_w / 2
    max_h = 5.625 - content_top - 0.1

    label_h = 0.3 if (left_label or right_label) else 0
    img_top = content_top + label_h

    for idx, (img_path, label) in enumerate([
        (left_image_path, left_label),
        (right_image_path, right_label),
    ]):
        x_offset = pad_side + idx * (half_w + gap)

        if label:
            _add_text_box(slide, Inches(x_offset), Inches(content_top),
                          Inches(half_w), Inches(0.25),
                          label, font_size=14, font_color=SC_BG, bold=True,
                          alignment=PP_ALIGN.CENTER)

        if img_path and os.path.exists(img_path):
            with PILImage.open(img_path) as img:
                img_w, img_h = img.size

            img_max_h = max_h - label_h
            img_aspect = img_w / img_h
            area_aspect = half_w / img_max_h

            if img_aspect > area_aspect:
                fit_w = half_w
                fit_h = half_w / img_aspect
            else:
                fit_h = img_max_h
                fit_w = img_max_h * img_aspect

            # Centre horizontally and vertically within half
            left = x_offset + (half_w - fit_w) / 2
            top = img_top + (img_max_h - fit_h) / 2

            slide.shapes.add_picture(
                img_path,
                Inches(left), Inches(top),
                width=Inches(fit_w), height=Inches(fit_h),
            )


# ──────────────────────────────────────────────────────────────────────
# Main report generator
# ──────────────────────────────────────────────────────────────────────

def generate_report(
        team_name: str,
        sections: list,
        output_path: str = "team_report.pptx",
        report_title: str = None,
        report_subtitle: str = None,
        template_path: str = None,
):
    """
    Generate a PPTX report from HTML visualizations or pre-saved images.
    HTML strings are automatically rendered to PNG via Chrome headless.

    Uses SkillCorner branded assets (backgrounds, logos) extracted from
    the template and stored in template_assets/.

    Parameters
    ----------
    team_name : str
        The team name for the report.
    sections : list of dict
        Each dict has:
          - type: "content" (default) or "divider"
          - title: str
          - html: str (HTML string — rendered to PNG automatically)
          - image_path: str (pre-saved PNG — takes priority over html)
          - subtitle: str (optional)
          - size: tuple (viewport w,h for HTML rendering, default 1400x900)
    output_path : str
        Path to save the output PPTX file.
    report_title : str, optional
        Custom title for the report.
    report_subtitle : str, optional
        Subtitle text on the title slide.
    template_path : str, optional
        Path to a SkillCorner PPTX template. Slide dimensions and theme
        are inherited. All existing slides are removed before adding new ones.
    """
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        # Remove all existing slides from the template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get('r:id')
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    _add_title_slide(prs, team_name, report_title, report_subtitle)

    tmp_dir = tempfile.mkdtemp()
    tmp_files = []
    section_counter = 0

    for i, section in enumerate(sections):
        slide_type = section.get("type", "content")

        if slide_type == "divider":
            section_counter += 1
            _add_section_divider(prs, section.get("title", ""),
                                 section_number=section_counter)
        elif slide_type == "dual":
            # Dual image slide (e.g. two radars side by side)
            def _resolve_image(key_prefix, idx):
                img_p = section.get(f"{key_prefix}_image_path", "")
                html_s = section.get(f"{key_prefix}_html", "")
                if (not img_p or not os.path.exists(img_p)) and html_s:
                    tp = os.path.join(tmp_dir, f"slide_{idx}_{key_prefix}.png")
                    sz = section.get(f"{key_prefix}_size", (700, 700))
                    _html_to_image(html_s, tp, size=sz)
                    tmp_files.append(tp)
                    return tp
                return img_p

            left_img = _resolve_image("left", i)
            right_img = _resolve_image("right", i)

            _add_dual_image_slide(
                prs,
                title=section.get("title", ""),
                left_image_path=left_img,
                right_image_path=right_img,
                left_label=section.get("left_label", None),
                right_label=section.get("right_label", None),
            )
        else:
            image_path = section.get("image_path", "")
            html_string = section.get("html", "")

            if (not image_path or not os.path.exists(image_path)) \
                    and html_string:
                tmp_path = os.path.join(tmp_dir, f"slide_{i}.png")
                size = section.get("size", (1400, 1100))
                _html_to_image(html_string, tmp_path, size=size)
                image_path = tmp_path
                tmp_files.append(tmp_path)

            _add_content_slide(
                prs,
                title=section.get("title", ""),
                image_path=image_path,
                subtitle=section.get("subtitle", None),
            )

    prs.save(output_path)

    # Clean up temp files
    for f in tmp_files:
        try:
            os.remove(f)
        except OSError:
            pass
    try:
        os.rmdir(tmp_dir)
    except OSError:
        pass

    return output_path
